#!/usr/bin/env python3
"""RecycleAI 프로젝트 발표용 PPTX 생성.

기본(HTML 동기화·애니메이션 캡처):
  python docs/presentation/build_presentation.py

텍스트-only PPT:
  python docs/presentation/build_presentation.py --text [--animate]
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "RecycleAI_Project_Deck.pptx"

# 브랜드
GREEN = RGBColor(0x2D, 0x5A, 0x27)
GREEN_LIGHT = RGBColor(0x10, 0xB9, 0x81)
DARK = RGBColor(0x1E, 0x29, 0x3B)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xFC)


def _blank(prs: Presentation):
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def _bar(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.05))
    box.fill.solid()
    box.fill.fore_color.rgb = GREEN
    box.line.fill.background()
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)


def _bullets(slide, items: list[str], top=1.35, left=0.7, width=12.0, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)


def _center_title(slide, title: str, sub: str | None = None, title_kr: str | None = None):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2))
    tf = t.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = GREEN
    if title_kr:
        kr = slide.shapes.add_textbox(Inches(0.8), Inches(3.35), Inches(11.7), Inches(0.6))
        kf = kr.text_frame
        kf.text = title_kr
        kp = kf.paragraphs[0]
        kp.alignment = PP_ALIGN.CENTER
        kp.font.size = Pt(24)
        kp.font.bold = True
        kp.font.color.rgb = GREEN
    if sub:
        sub_top = 4.15 if title_kr else 3.6
        s = slide.shapes.add_textbox(Inches(1.2), Inches(sub_top), Inches(10.9), Inches(1.5))
        sf = s.text_frame
        sf.text = sub
        sp = sf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(18)
        sp.font.color.rgb = GRAY


def _table(slide, headers: list[str], rows: list[list[str]], top=1.5):
    n_rows, n_cols = len(rows) + 1, len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(top), Inches(12.3), Inches(0.55 * n_rows))
    tbl = shape.table
    for c, h in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = DARK
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG


def _image_or_placeholder(slide, path: Path | None, left, top, w, h, caption: str, placeholder: str):
    if path and path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(w))
    else:
        ph = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        ph.line.color.rgb = GRAY
        tf = ph.text_frame
        tf.text = placeholder
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cap = slide.shapes.add_textbox(Inches(left), Inches(top + h + 0.05), Inches(w), Inches(0.45))
    cap.text_frame.text = caption
    cap.text_frame.paragraphs[0].font.size = Pt(11)
    cap.text_frame.paragraphs[0].font.color.rgb = GRAY
    cap.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def _video_or_placeholder(
    slide,
    path: Path | None,
    left: float,
    top: float,
    w: float,
    h: float,
    placeholder: str,
) -> None:
    if path and path.exists():
        slide.shapes.add_movie(
            str(path),
            Inches(left),
            Inches(top),
            width=Inches(w),
            height=Inches(h),
            mime_type="video/mp4",
        )
    else:
        ph = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
        ph.fill.solid()
        ph.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        ph.line.color.rgb = GRAY
        tf = ph.text_frame
        tf.text = placeholder
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = GRAY
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def build(animated: bool = False) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 타이틀
    s = _blank(prs)
    _center_title(
        s,
        "RecycleAI",
        "카메라만 비추면, 우리 동네 기준 분리배출을 안내합니다\n바이브코딩 프로젝트 발표 | 2026",
        title_kr="리사이클AI",
    )
    foot = s.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5))
    foot.text_frame.text = "발표자: _________  |  연락: _________"
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    foot.text_frame.paragraphs[0].font.size = Pt(12)
    foot.text_frame.paragraphs[0].font.color.rgb = GRAY

    # 2 Problem
    s = _blank(prs)
    _bar(s, "왜 만들었나 — Problem", "이사·출장·여행으로 ‘동네 규칙’이 바뀔 때")
    _bullets(
        s,
        [
            "페르소나: 35년 강남 생활 후 일산 이사한 60대 — 분리수거 경고 스티커를 두 번 받음",
            "고무장갑·즉석밥 용기 등 ‘예전 동네 습관’이 오히려 실수를 만듦",
            "지자체마다 분류가 반대인 경우도 있음 → 과태료·수거 거부·수치심",
            "하루 6~12회, 분리수거일엔 15~25개 — 의심 케이스만 1~3분 × 주당 10~15회",
            "기존 대안: PDF(글자 작음), 이웃·경비실 문의(부담), 검색 앱(타 지역 기준)",
        ],
        size=19,
    )

    # 3 Value
    s = _blank(prs)
    _bar(s, "사용자가 되찾는 가치")
    _bullets(
        s,
        [
            "자존심 회복 — “모르는 게 아니라 동네가 다르다”는 안심",
            "과태료·경고 리스크 감소 (건당 수만 원대 가능)",
            "분리수거 시간 30분 → 10분 수준 단축 (가정)",
            "이웃·경비실 반복 질문 없이 스스로 판단",
            "이사·출장 시에도 GPS 기준으로 즉시 현지 규칙 적용",
        ],
    )

    # 4 Solution
    s = _blank(prs)
    _bar(s, "Solution — RecycleAI 한 줄")
    _center_title(s, "버리려는 물건을 카메라로 비추면", "AI 인식 + 우리 동네 SQLite 규칙 → 한 화면 카드에 배출법·요일·문의처")
    steps = s.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.3), Inches(2.2))
    tf = steps.text_frame
    for i, t in enumerate(
        [
            "① GPS로 지자체 확인",
            "② ML Kit 초록 박스 / 손가락 주황 박스로 물건 지정",
            "③ 로컬 DB + (필요 시) Gemini → 픽토그램·배출 요일·E-순환 안내",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

    # 5 개발 프로세스
    s = _blank(prs)
    _bar(s, "일반적인 앱 개발 프로세스 (본 프로젝트 적용)")
    _table(
        s,
        ["단계", "우리가 한 일", "산출물"],
        [
            ["1. 문제·페르소나", "박정호 케이스 구체화", "problem.md"],
            ["2. 기획·스펙", "화면별 P0 기능 정의", "PRD.md, vision.md"],
            ["3. UX·디자인", "글래스모피즘·네온 박스·핀치 줌", "ui/ 프로토타입, Tokens.kt"],
            ["4. 데이터", "공공 품목·조례·MOIS 크롤", "wasteguide.sqlite3"],
            ["5. 구현", "Compose + ML Kit + Gemini", "Android APK"],
            ["6. 검증·출시", "실기기 테스트 → 스토어", "Play Console, 홍보"],
        ],
        top=1.25,
    )

    # 6 AI 개발 방식
    s = _blank(prs)
    _bar(s, "개발 방식 — 문서 먼저, AI와 함께 구현")
    _bullets(
        s,
        [
            "problem / PRD / TRD를 먼저 써서 AI·개발자가 같은 ‘제품 언어’를 쓰게 함",
            "카드 문구는 공공 DB(SQLite)만 사용 — AI 답변을 그대로 노출하지 않음 (신뢰)",
            "로컬 730품목 1차 매칭 → 애매할 때만 Gemini (비용·발열 절감)",
            "오프라인 우선: 지하철·와이파이 없어도 조례·배출 요일 조회 가능",
        ],
    )

    # 7 아키텍처
    s = _blank(prs)
    _bar(s, "기술 구조 (쉬운 말로)")
    _bullets(
        s,
        [
            "카메라: Android CameraX + Google ML Kit (가벼운 객체 감지)",
            "데이터: 앱 안 SQLite — 품목 730+, 지자체 조례, MOIS 배출 일정",
            "AI: Gemini는 ‘사진→한글 키워드’만, 안내 문장은 DB에서 조합",
            "위치: GPS + Geocoder → 해당 구·시 규칙 카드",
            "폐가전: E-순환 무상 수거 안내·전화 연결",
        ],
    )

    # 8 경쟁 앱 소개
    s = _blank(prs)
    _bar(s, "비교 대상 — 「내 손안의 분리배출」", "환경부·한국환경공단 등 공식 앱 | Play 10만+ 다운로드")
    _bullets(
        s,
        [
            "강점: 공신력, 재질 8분류, 품목 검색, 전국 기본 가이드",
            "한계(우리 관찰): 텍스트·카테고리 중심, 실시간 카메라 인식 없음",
            "지역: 전국 ‘표준’ 안내 — 거주지 GPS·조례·배출 요일 일체형은 아님",
            "네트워크: 온라인 환경 권장(보도·리뷰 기준)",
        ],
        top=1.2,
        size=18,
    )
    _image_or_placeholder(
        s,
        ASSETS / "competitor_home.png",
        8.2,
        1.35,
        4.6,
        4.8,
        "공식 앱 첫 화면 (대한민국 정책브리핑 캡처)",
        "",
    )

    # 9 비교 표
    s = _blank(prs)
    _bar(s, "기능 비교 한눈에")
    _table(
        s,
        ["항목", "내 손안의 분리배출", "RecycleAI"],
        [
            ["입력 방식", "검색·카테고리 선택", "카메라 실시간 + 드래그 영역"],
            ["지역 맞춤", "전국 공통 기준", "GPS → 해당 지자체 조례·배출일"],
            ["오프라인", "제한적(온라인 권장)", "품목·조례·일정 로컬 DB"],
            ["시각 안내", "정적 사진·텍스트", "네온 박스 + 픽토그램 카드"],
            ["접근성", "일반 UI", "핀치 줌·큰 카드·전화 바로걸기"],
            ["폐가전", "별도 안내", "스캔 시 E-순환·문의처 연동"],
        ],
        top=1.2,
    )

    # 10 경쟁 스크린샷 1
    s = _blank(prs)
    _bar(s, "경쟁 앱 화면 — 카테고리·품목 안내")
    _image_or_placeholder(s, ASSETS / "competitor_categories.png", 0.6, 1.3, 5.8, 5.2, "재질별 분류 (정책브리핑)", "")
    _image_or_placeholder(
        s,
        ASSETS / "recycleai_scan.png",
        6.9,
        1.3,
        5.8,
        5.2,
        "RecycleAI — 카메라 스캔·초록 박스",
        "[캡처 삽입]\n실기기 스크린샷",
    )

    # 11 우리만의 장점
    s = _blank(prs)
    _bar(s, "RecycleAI가 어필할 5가지")
    _bullets(
        s,
        [
            "「지금 이 동네」 — 이사·출장 사용자에게 가장 큰 차별점",
            "손에 든 물건 그대로 — 검색어를 고민할 필요 없음",
            "한 화면 원스톱 — 분류·세척·라벨·배출 요일·지자체 전화",
            "신뢰 설계 — 공공 DB 문구만 표시, AI는 보조 매칭만",
            "시니어 UX — 핀치 줌, 네온 피드백, 다시 스캔",
        ],
        size=22,
    )

    # 12 데모 시나리오 + 시연 영상
    s = _blank(prs)
    _bar(s, "데모 시나리오 (발표 시 연출)")
    _bullets(
        s,
        [
            "앱 실행 → 상단 「고양시 일산서구」 확인",
            "마우스 터치 → 초록 박스 → 카드에 배출법·요일",
            "플라스틱 컵 스캔 → 비우기·헹구기 픽토그램",
            "여러 물건 겹침 → 주황 드래그 박스로 한 개만 분석",
            "폐가전 스캔 → E-순환 무상 수거·1599-0903",
        ],
        left=1.8,
        width=5.4,
        size=18,
    )
    _video_or_placeholder(
        s,
        ASSETS / "recycle_demo.mp4",
        7.35,
        1.45,
        4.2,
        5.55,
        "[시연 영상]\nassets/recycle_demo.mp4",
    )

    # 14 현황
    s = _blank(prs)
    _bar(s, "현재 상태 (2026)")
    _bullets(
        s,
        [
            "Android Compose 앱 — Production Ready 수준 문서·DB 파이프라인",
            "품목·조례: wasteguide 기반 SQLite 번들",
            "배출 일정: MOIS 데이터 일부 지자체 적재(확장 중)",
            "다음: Play 스토어 등록, 실사용 피드백, iOS·전국 커버리지",
        ],
    )

    # 15 Play Store
    s = _blank(prs)
    _bar(s, "Google Play 배포 로드맵")
    _bullets(
        s,
        [
            "개발자 계정 등록(일회성 등록비) · 앱 서명 키 안전 보관",
            "스토어 등록정보: 앱명 RecycleAI, 짧은 설명·스크린샷 6~8장·정책 URL",
            "개인정보 처리방침·카메라·위치 권한 고지 (Play 정책 필수)",
            "내부 테스트 → 클로즈드 테스트(지인 20명) → 프로덕션 단계 출시",
            "ASO 키워드: 분리수거, 분리배출, 재활용, 이사, 지역, 카메라",
        ],
        size=18,
    )

    # 16 SNS 홍보
    s = _blank(prs)
    _bar(s, "SNS·콘텐츠 홍보")
    _bullets(
        s,
        [
            "숏폼(릴스·틱톡·쇼츠): 「고무장갑 강남 vs 일산」 15초 비교",
            "카카오 오픈채팅·아파트 커뮤니티: 이사 시즌(2~3월, 9월) 집중",
            "블로그·네이버: 「분리수거 경고 스티커」 검색어 SEO",
            "지자체·환경 박람회·대학 EC 동아리 시연 부스",
            "인플루언서: 육아·생활·시니어 IT 유튜버 체험 영상",
        ],
        size=18,
    )

    # 17 수익화
    s = _blank(prs)
    _bar(s, "수익화 방향 (단계별)")
    _table(
        s,
        ["단계", "모델", "설명"],
        [
            ["1단계", "무료 + 브랜드", "사용자·리뷰 확보, 공공 미션 이미지"],
            ["2단계", "프리미엄", "광고 제거, 이사 모드(출신지↔현재지 비교), 가족 공유"],
            ["3단계", "B2B·B2G", "아파트·지자체 화이트라벨, 데이터 리포트(익명 통계)"],
            ["4단계", "제휴", "대형폐기물·리사이클 업체 연계(추천 수수료, 신중히)"],
        ],
        top=1.25,
    )
    note = s.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(12), Inches(0.8))
    note.text_frame.text = "※ 공공 데이터·환경 교육 앱 특성상, 과도한 광고보다 신뢰·파트너십 우선 권장"
    note.text_frame.paragraphs[0].font.size = Pt(13)
    note.text_frame.paragraphs[0].font.color.rgb = GRAY

    # 18 로드맵
    s = _blank(prs)
    _bar(s, "12개월 로드맵")
    _bullets(
        s,
        [
            "Q1: Play 출시, MOIS·조례 커버리지 확대, 스토어 리뷰 대응",
            "Q2: 이사 모드·품목 Clarification 고도화, iOS 검토",
            "Q3: 지자체·아파트 파일럿, B2G 제안서",
            "Q4: 전국 커버리지 목표, 프리미엄·제휴 파일럿",
        ],
        size=20,
    )

    # 19 프로젝트가 주는 메시지
    s = _blank(prs)
    _bar(s, "이 프로젝트가 주는 메시지")
    _bullets(
        s,
        [
            "문제는 ‘분류 지식’이 아니라 ‘지역 이동’ — 인구 이동·고령화와 맞물림",
            "공공 데이터 + 모바일 AI의 조합 — 학생도 만들 수 있는 파이프라인",
            "공식 앱과 공존·차별 — 검색 앱을 대체가 아닌 ‘현장 카메라 레이어’",
            "확장 아이디어: 지자체·학교 협력, 교육·ESG, 해외 교포·유학생 활용",
        ],
        size=20,
    )

    # 20 Supabase 정정
    s = _blank(prs)
    _bar(s, "중간발표 정정 — Supabase는 런타임 검색 DB가 아닙니다", "수정·추가 하이라이트")
    _table(
        s,
        ["구분", "중간발표에서 생긴 오해", "현재 발표 문서의 정정 표현"],
        [
            ["앱 실행 중", "Supabase에 연결해 실시간 사물 검색", "Supabase·PostgREST·백엔드 없음. 앱은 로컬 SQLite를 조회"],
            ["AI 역할", "AI가 안내 문장을 직접 생성", "Gemini는 사진을 한글 키워드로 보조 변환. 안내문은 DB 문구만 표시"],
            ["DB 갱신", "클라우드 DB를 계속 읽음", "행안부 DB가 730여 개에서 1500개로 확장되면, 배포 전 로컬 DB 업데이트 작업에만 임시 활용 가능"],
            ["지역 안내", "전국 공통 안내", "품목·조례·MOIS 일정·문의처를 통합 SQLite로 묶고 GPS 지역 기준으로 분리 안내"],
        ],
        top=1.25,
    )
    note = s.shapes.add_textbox(Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.6))
    note.text_frame.text = "발표 멘트: “Supabase는 현재 앱 런타임 검색에 쓰지 않습니다. 향후 공공 DB 확장 시, 앱 안 SQLite를 재생성하는 임시 데이터 작업 도구로만 검토합니다.”"
    note.text_frame.paragraphs[0].font.size = Pt(13)
    note.text_frame.paragraphs[0].font.color.rgb = GREEN

    # 21 Eval 테스트 케이스
    s = _blank(prs)
    _bar(s, "Eval 테스트 케이스 — 낮은 점수 항목은 이렇게 고쳤습니다")
    _table(
        s,
        ["평가 항목", "테스트 기준", "수정 전", "수정 후"],
        [
            ["핀치 줌", "확대/축소 후 하단 고지까지 스크롤", "확대 시 하단 잘림", "layout 기반 줌 컬럼으로 전체 스크롤"],
            ["Supabase 배제", "비행기 모드에서도 규칙 카드 표시", "클라우드 검색처럼 설명됨", "SQLite 번들 + Gemini 보조로 명확화"],
            ["DB 통합", "품목·지역·MOIS·문의처가 한 카드에 결합", "지역별 안내 분산", "통합 SQLite에서 지역 기준 매칭"],
            ["반복 안내문", "동일 주의문은 1회만 노출", "세척·라벨 안내 반복", "공통 가이드를 1회 카드로 단순화"],
            ["오터치 방어", "작은 드래그·빈 화면 탭에서 복구", "작은 박스 후 UX 막힘", "32px 이하 취소, 배경 탭 해제"],
        ],
        top=1.15,
    )

    # 22 앱 실행 흐름
    s = _blank(prs)
    _bar(s, "앱 실행 흐름 — 사용자는 쉽게, 내부는 신뢰 데이터로")
    process_steps = [
        ("APP", "1. 아이콘 클릭", "CameraX 카메라\nGPS/Geocoder 준비"),
        ("CAM", "2. 물건 비추기", "ML Kit이 후보 물체를\n초록 박스로 추적"),
        ("AI", "3. 애매할 때 Gemini", "이미지를 한글 키워드로\n보조 변환"),
        ("TAP", "4. 사물 클릭", "키워드와 품목 규칙\nDB 매칭"),
        ("DB", "5. 지역 DB 결합", "조례·MOIS·문의처를\nSQLite에서 조회"),
        ("CARD", "6. 결과 카드", "분류·배출 요일·전화\n한 화면 표시"),
    ]
    for idx, (icon, title, body) in enumerate(process_steps):
        if idx < 3:
            col = idx
            row = 0
        else:
            col = 5 - idx
            row = 1
        left = 0.75 + col * 4.15
        top = 1.35 + row * 2.45
        card = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(3.6), Inches(1.95))
        card.fill.solid()
        card.fill.fore_color.rgb = BG
        card.line.color.rgb = GREEN_LIGHT
        badge = s.shapes.add_shape(1, Inches(left + 0.18), Inches(top + 0.18), Inches(0.75), Inches(0.45))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        badge.text_frame.text = icon
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = WHITE
        title_box = s.shapes.add_textbox(Inches(left + 1.05), Inches(top + 0.18), Inches(2.25), Inches(0.35))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(14)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = DARK
        body_box = s.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.85), Inches(3.1), Inches(0.8))
        body_box.text_frame.text = body
        for p in body_box.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = GRAY
            p.alignment = PP_ALIGN.CENTER
    for left, top, mark in [
        (4.42, 2.12, "→"),
        (8.57, 2.12, "→"),
        (10.62, 3.28, "↓"),
        (8.57, 4.57, "←"),
        (4.42, 4.57, "←"),
    ]:
        arrow = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(0.36), Inches(0.36))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
        arrow.line.color.rgb = GREEN_LIGHT
        arrow.text_frame.text = mark
        arrow.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        arrow.text_frame.paragraphs[0].font.size = Pt(13)
        arrow.text_frame.paragraphs[0].font.bold = True
        arrow.text_frame.paragraphs[0].font.color.rgb = GREEN

    # 23 초기 무비용 마케팅
    s = _blank(prs)
    _bar(s, "초기 비용 0원 마케팅 실행법", "돈보다 중요한 것은 ‘누가 왜 써야 하는지’를 정확히 보여주는 것")
    _table(
        s,
        ["방법", "바로 할 일", "기대 효과"],
        [
            ["스토어 ASO", "앱 제목·설명에 분리수거, 이사, 지역, 카메라 키워드 반영", "검색 유입 확보"],
            ["지인 20명 테스트", "가족·학우·아파트 이웃에게 설치 요청 후 리뷰 받기", "초기 신뢰와 오류 발견"],
            ["짧은 시연 영상", "고무장갑/종이컵/폐가전 15초 전후 비교 영상 제작", "앱 가치를 한눈에 전달"],
            ["커뮤니티 공유", "아파트 카톡방·맘카페·지역 생활 게시판에 도움형 글 작성", "실사용자 반응 수집"],
            ["블로그 글 3개", "‘일산 분리수거’, ‘이사 후 분리배출’ 같은 검색형 글 작성", "무료 검색 노출"],
            ["피드백 루프", "불편 의견을 관리 페이지·스토어 업데이트 노트에 반영", "사용자가 개선을 체감"],
        ],
        top=1.15,
    )
    tip = s.shapes.add_textbox(Inches(0.7), Inches(6.25), Inches(12.0), Inches(0.55))
    tip.text_frame.text = "초기 목표: 광고비 지출보다 ‘실제 써본 사람의 후기 20개’와 ‘짧은 시연 콘텐츠 3개’를 먼저 확보합니다."
    tip.text_frame.paragraphs[0].font.size = Pt(13)
    tip.text_frame.paragraphs[0].font.bold = True
    tip.text_frame.paragraphs[0].font.color.rgb = GREEN

    # 24 배포 후 광고 및 관리
    s = _blank(prs)
    _bar(s, "배포 후 광고 및 관리", "운영자가 코드·Remote Config로 노출과 사용량을 제어")
    _image_or_placeholder(
        s,
        ASSETS / "recycleAI_02.png",
        0.55,
        1.25,
        3.1,
        5.1,
        "사용자 화면: 스캔·결과 카드",
        "[앱 캡처]\n스캔/결과 화면",
    )

    admin = s.shapes.add_shape(1, Inches(4.05), Inches(1.25), Inches(4.15), Inches(5.1))
    admin.fill.solid()
    admin.fill.fore_color.rgb = RGBColor(0xEC, 0xFD, 0xF5)
    admin.line.color.rgb = GREEN_LIGHT
    admin.text_frame.text = "관리 설정"
    admin.text_frame.paragraphs[0].font.size = Pt(20)
    admin.text_frame.paragraphs[0].font.bold = True
    admin.text_frame.paragraphs[0].font.color.rgb = GREEN

    settings = [
        ("광고 배너", "SHOW_AD_BANNER\nfalse ↔ true"),
        ("AdMob ID", "테스트 ID → 실제 광고 단위 ID"),
        ("횟수제한", "Remote Config\nlimit_enabled"),
        ("무료 횟수", "daily_scan_limit = 5"),
    ]
    for i, (title, body) in enumerate(settings):
        top = 2.0 + i * 0.88
        card = s.shapes.add_shape(1, Inches(4.35), Inches(top), Inches(3.55), Inches(0.68))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xBB, 0xF7, 0xD0)
        label = s.shapes.add_textbox(Inches(4.55), Inches(top + 0.1), Inches(1.2), Inches(0.25))
        label.text_frame.text = title
        label.text_frame.paragraphs[0].font.size = Pt(11)
        label.text_frame.paragraphs[0].font.bold = True
        label.text_frame.paragraphs[0].font.color.rgb = DARK
        value = s.shapes.add_textbox(Inches(5.75), Inches(top + 0.08), Inches(1.95), Inches(0.45))
        value.text_frame.text = body
        for p in value.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.color.rgb = GRAY

    _bullets(
        s,
        [
            "배너 광고: 결과 카드 영역에 AdMob 배너를 붙이고, 플래그 하나로 표시/숨김 전환",
            "횟수제한: Firebase Remote Config에서 on/off와 일일 무료 스캔 수를 원격 조정",
            "한도 초과 시: 보상형 광고 안내 → 시청 완료 후 로컬 카운트 충전 → 스캔 재개",
            "출시 전 체크: 실제 광고 ID, 개인정보 처리방침, 광고 SDK 정책 검수",
        ],
        top=1.32,
        left=8.45,
        width=4.3,
        size=14,
    )
    note = s.shapes.add_textbox(Inches(4.1), Inches(5.85), Inches(8.65), Inches(0.62))
    note.text_frame.text = "운영 문서: docs/admin_operations_guide.md"
    note.text_frame.paragraphs[0].font.size = Pt(13)
    note.text_frame.paragraphs[0].font.bold = True
    note.text_frame.paragraphs[0].font.color.rgb = GREEN

    # 25 무료 로컬 TTS 나레이션
    s = _blank(prs)
    _bar(s, "시연 영상 나레이션 — 내 PC의 무료 TTS로 제작", "유료 성우·구독 없이, 로컬에서 음성 합성")
    _bullets(
        s,
        [
            "사용 모델: Supertone 'supertonic-3' (Hugging Face 공개 모델)",
            "특징: 내 컴퓨터에 설치해 무료로 사용, 한국어 지원, GPU 없이 CPU로 동작",
            "동작: 첫 1회만 모델 다운로드 → 이후 오프라인으로 WAV 합성",
            "방법: pip install supertonic → 프리셋 목소리(M1~M5 / F1~F5) 선택 → 문장 입력 → WAV 저장",
            "결과물: 시연 영상에 깔린 나레이션 음성을 이 방식으로 직접 생성",
        ],
        top=1.3,
        size=18,
    )
    link = s.shapes.add_textbox(Inches(0.7), Inches(4.25), Inches(11.9), Inches(0.4))
    link.text_frame.text = "모델 주소: https://huggingface.co/Supertone/supertonic-3"
    link.text_frame.paragraphs[0].font.size = Pt(14)
    link.text_frame.paragraphs[0].font.bold = True
    link.text_frame.paragraphs[0].font.color.rgb = GREEN_LIGHT
    code = s.shapes.add_shape(1, Inches(0.7), Inches(4.75), Inches(11.9), Inches(2.05))
    code.fill.solid()
    code.fill.fore_color.rgb = DARK
    code.line.fill.background()
    ctf = code.text_frame
    ctf.word_wrap = True
    code_lines = [
        "# 설치 (한 번만)",
        "pip install supertonic onnxruntime",
        "",
        "# 나레이션 만들기 (한국어 프리셋 M1)",
        "python make_narration.py --file 원고.txt --voice M1 --out narration.wav",
    ]
    for i, line in enumerate(code_lines):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(0x86, 0xEF, 0xAC) if line.startswith("#") else WHITE

    # 26 바이브코딩 사용 툴·모델 정리 (초보자용)
    s = _blank(prs)
    _bar(s, "이 앱은 무엇으로 만들었나 — 바이브코딩 도구 정리", "초보자도 따라 할 수 있도록, 쓴 도구·모델을 한눈에")
    _table(
        s,
        ["분류", "사용한 도구 / 모델", "한 일 (쉽게)"],
        [
            ["AI 코딩", "Cursor (AI 코드 에디터)", "대화하듯 코드 작성·수정"],
            ["앱 개발", "Kotlin · Jetpack Compose", "안드로이드 화면·기능 구현"],
            ["카메라 인식", "CameraX · Google ML Kit", "사물 비추면 박스로 인식"],
            ["AI 보조 인식", "Google Gemini API", "사진 → 한글 키워드 변환"],
            ["데이터", "SQLite + Python 크롤링", "품목·지역 규칙을 앱 안에 저장"],
            ["광고·설정", "Google AdMob · Firebase Remote Config", "배너 광고·무료 횟수 원격 조정"],
            ["나레이션", "Supertonic 3 (무료 로컬 TTS)", "시연 영상 음성 합성"],
            ["발표 자료", "python-pptx", "이 슬라이드 자동 생성"],
        ],
        top=1.2,
    )
    tip = s.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.6))
    tip.text_frame.text = "핵심: 대부분 무료/공개 도구로 구성 — 유료 서버 없이 'AI와 대화하며' 앱·영상·발표까지 완성"
    tip.text_frame.paragraphs[0].font.size = Pt(13)
    tip.text_frame.paragraphs[0].font.bold = True
    tip.text_frame.paragraphs[0].font.color.rgb = GREEN

    # 27 SNS 홍보 영상 (UGC, Q&A 직전)
    s = _blank(prs)
    _bar(s, "SNS 홍보 영상 (15초)", "시연(12번)과 별도 · ▶ 재생 후 음소거 해제 권장")
    _video_or_placeholder(
        s,
        ASSETS / "recycle_ad.mp4",
        4.0,
        1.35,
        5.2,
        5.7,
        "[홍보 영상]\nassets/recycle_ad.mp4\n9:16 UGC",
    )
    hint = s.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12.0), Inches(0.55))
    hint.text_frame.text = "발표: 이 슬라이드에서 재생 → Q&A(다음)로 이동 · 파일 직접 실행: assets/recycle_ad.mp4"
    hint.text_frame.paragraphs[0].font.size = Pt(12)
    hint.text_frame.paragraphs[0].font.color.rgb = GRAY

    # 28 Q&A
    s = _blank(prs)
    _center_title(s, "감사합니다", "바이브코딩을 배우는 우리 모두에게")
    qa = s.shapes.add_textbox(Inches(1.5), Inches(4.15), Inches(10.3), Inches(2.3))
    tf = qa.text_frame
    for i, t in enumerate(
        [
            "좋은 개발자는 정답을 외우는 사람이 아니라,",
            "문제를 끝까지 관찰하고 더 나은 질문을 던지는 사람입니다.",
            "오늘의 작은 실험이 내일의 진짜 서비스가 됩니다.",
        ]
    ):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18 if i < 2 else 16)
        p.font.bold = i < 2
        p.font.color.rgb = DARK if i < 2 else GRAY

    if animated:
        from pptx_effects import set_slide_fade_transition

        for i, slide in enumerate(prs.slides):
            if i > 0:
                set_slide_fade_transition(slide)

    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    import sys

    if "--text" in sys.argv:
        path = build(animated="--animate" in sys.argv)
        print(f"생성 완료 (텍스트): {path}")
    else:
        from sync_html_to_pptx import main as sync_main

        sync_main()
