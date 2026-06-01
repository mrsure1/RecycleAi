#!/usr/bin/env python3
"""
HTML 발표(docs/presentation/web/index.html)를 기준으로 PPTX 생성.

- visual (기본): Playwright로 각 슬라이드 CSS 애니메이션 재생 후 캡처 → 픽셀 단위 동기화
- text: python-pptx 텍스트 레이아웃 + 등장 애니메이션 (기존 build 로직)

실행:
  python docs/presentation/sync_html_to_pptx.py
  python docs/presentation/sync_html_to_pptx.py --mode text
"""

from __future__ import annotations

import argparse
import http.server
import socket
import threading
import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from pptx_effects import set_slide_fade_transition

ROOT = Path(__file__).resolve().parent
WEB = ROOT / "web"
ASSETS = ROOT / "assets"
OUT_VISUAL = ROOT / "RecycleAI_Project_Deck.pptx"
OUT_EDITABLE = ROOT / "RecycleAI_Project_Deck_editable.pptx"
OUT = OUT_VISUAL
CAPTURE_DIR = ROOT / "captures"
HTML_URL_PATH = "/web/index.html"

SLIDE_W_PX = 1280
SLIDE_H_PX = 720


def _free_port() -> int:
    with socket.socket(socket.AF_INET, socket.SOCK_STREAM) as s:
        s.bind(("127.0.0.1", 0))
        return s.getsockname()[1]


def _start_server(port: int) -> http.server.ThreadingHTTPServer:
    handler = http.server.SimpleHTTPRequestHandler

    class Handler(handler):
        def __init__(self, *args, **kwargs):
            # presentation/ 루트 → /web/index.html, /assets/* 모두 제공
            super().__init__(*args, directory=str(ROOT), **kwargs)

    server = http.server.ThreadingHTTPServer(("127.0.0.1", port), Handler)
    thread = threading.Thread(target=server.serve_forever, daemon=True)
    thread.start()
    return server


def capture_slides_from_html() -> list[Path]:
    """Playwright로 활성 슬라이드 스크린샷 (애니메이션 포함)."""
    from playwright.sync_api import sync_playwright

    CAPTURE_DIR.mkdir(parents=True, exist_ok=True)
    port = _free_port()
    server = _start_server(port)
    base = f"http://127.0.0.1:{port}{HTML_URL_PATH}"

    paths: list[Path] = []
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(
                viewport={"width": SLIDE_W_PX, "height": SLIDE_H_PX},
                device_scale_factor=1,
            )
            page.goto(base, wait_until="networkidle", timeout=120_000)
            page.wait_for_function(
                "() => typeof window.goToSlide === 'function'",
                timeout=30_000,
            )
            # 1페이지: 정적 프로 타이틀 (GIF·시네마틱 없음)
            page.evaluate("() => document.body.classList.add('ppt-capture')")
            total = page.evaluate(
                "() => document.querySelectorAll('#presentation > .slide').length"
            )
            for i in range(total):
                page.evaluate(f"window.goToSlide({i})")
                wait_ms = 600 if i == 0 else (800 if i == total - 1 else 2200)
                page.wait_for_timeout(wait_ms)
                out = CAPTURE_DIR / f"slide_{i + 1:02d}.png"
                if i == total - 1:
                    page.locator("#presentation").screenshot(path=str(out), type="png")
                else:
                    page.locator("#presentation > .slide.active").screenshot(
                        path=str(out), type="png"
                    )
                paths.append(out)
            browser.close()
    finally:
        server.shutdown()
    return paths


def build_visual_pptx(images: list[Path], *, embed_video: bool = True) -> Path:
    """캡처 PNG를 전체 슬라이드로 삽입 + Fade 전환."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]

    for idx, img in enumerate(images):
        if not img.is_file():
            continue
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_picture(
            str(img),
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )
        if idx > 0:
            set_slide_fade_transition(slide)

    prs.save(OUT)

    if embed_video:
        from pptx_embed_videos import embed_videos

        embed_videos(OUT)

    return OUT


def build_text_pptx() -> Path:
    """편집 가능 텍스트·표·이미지·영상 슬라이드 (python-pptx)."""
    from build_presentation import build as build_text

    return build_text(animated=True, output=OUT_EDITABLE)


def main() -> None:
    parser = argparse.ArgumentParser(description="HTML 발표 → PPTX")
    parser.add_argument(
        "--mode",
        choices=("visual", "text"),
        default="visual",
        help="visual=Playwright 캡처(애니메이션 동일), text=편집 가능 텍스트",
    )
    args = parser.parse_args()

    if args.mode == "text":
        path = build_text_pptx()
        print(f"편집 가능 PPT 생성: {path}")
        print("PowerPoint에서 글자·표·도형을 직접 수정할 수 있습니다.")
        print("시연·광고 영상은 12·25번 슬라이드에 MP4가 삽입됩니다.")
        return

    print("HTML 슬라이드 캡처 중 (Playwright)…")
    t0 = time.time()
    images = capture_slides_from_html()
    print(f"캡처 {len(images)}장 ({time.time() - t0:.1f}s)")
    path = build_visual_pptx(images)
    print(f"시각 동기화 PPT: {path}")
    print("슬라이드쇼에서 전환 효과는 Fade(0.7s)입니다.")
    print("1페이지: 정적 프로 타이틀 (네이비·틸 악센트)")
    print("12·25페이지: recycle_demo / recycle_ad (슬라이드쇼에서 ▶ 재생)")


if __name__ == "__main__":
    main()
