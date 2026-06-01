#!/usr/bin/env python3
"""12·25페이지 HTML 스타일 재디자인 + MP4 슬라이드쇼 재생."""

from __future__ import annotations

import re
import subprocess
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
HF_TTS_MODEL_URL = "https://huggingface.co/Supertone/supertonic-3"
HIGGSFIELD_URL = "https://higgsfield.ai"

# HTML :root 토큰
C_BG = RGBColor(0x0B, 0x11, 0x20)
C_SURFACE = RGBColor(0x11, 0x18, 0x27)
C_GREEN = RGBColor(0x2D, 0x5A, 0x27)
C_GREEN_LT = RGBColor(0x10, 0xB9, 0x81)
C_TEXT = RGBColor(0xF1, 0xF5, 0xF9)
C_TEXT2 = RGBColor(0x94, 0xA3, 0xB8)
C_MUTED = RGBColor(0x64, 0x74, 0x8B)
C_BORDER = RGBColor(0x33, 0x41, 0x55)

SLIDE_W = 13.333
SLIDE_H = 7.5

# 1280×720 HTML → 인치 (phone-frame 220×440, screen inset 8/20)
PHONE_W = 2.29
PHONE_H = 4.58
PHONE_LEFT_DEMO = 7.05
PHONE_LEFT_AD = (SLIDE_W - PHONE_W) / 2
PHONE_TOP = 1.72
SCREEN_INSET_X = 0.08
SCREEN_INSET_Y = 0.21
SCREEN_W = PHONE_W - SCREEN_INSET_X * 2
SCREEN_H = PHONE_H - SCREEN_INSET_Y * 2


@dataclass(frozen=True)
class TimelineStep:
    title: str
    desc: str


DEMO_STEPS = (
    TimelineStep("① 앱 실행", "상단에 「고양시 일산서구」 자동 표시 확인"),
    TimelineStep("② 물건 터치", "초록 박스 → 카드에 배출법·요일 표시"),
    TimelineStep("③ 플라스틱 컵 스캔", "비우기·헹구기 픽토그램 카드"),
    TimelineStep("④ 드래그 박스", "여러 물건 겹침 → 주황 드래그로 한 개만 분석"),
    TimelineStep("⑤ 폐가전 스캔", "E-순환 무상 수거 · 1599-0903 전화 연결"),
)


def _probe_video_size(path: Path) -> tuple[int, int]:
    ff = __import__("imageio_ffmpeg", fromlist=["get_ffmpeg_exe"]).get_ffmpeg_exe()
    proc = subprocess.run([ff, "-i", str(path)], capture_output=True, text=True, check=False)
    m = re.search(r"Video:.*? (\d{2,5})x(\d{2,5})", proc.stderr or "")
    if not m:
        raise RuntimeError(f"해상도 파싱 실패: {path}")
    return int(m.group(1)), int(m.group(2))


def _poster_for_video(video: Path) -> Path:
    """add_movie 포스터 없으면 PowerPoint가 스피커 아이콘을 늘려 표시함."""
    curated = ASSETS / "recycleAI_01.png"
    if video.name == "recycle_ad.mp4" and curated.is_file():
        return curated

    cache_dir = ASSETS / ".posters"
    cache_dir.mkdir(parents=True, exist_ok=True)
    poster = cache_dir / f"{video.stem}_frame0.jpg"
    if poster.is_file() and poster.stat().st_mtime >= video.stat().st_mtime:
        return poster

    ff = __import__("imageio_ffmpeg", fromlist=["get_ffmpeg_exe"]).get_ffmpeg_exe()
    subprocess.run(
        [ff, "-y", "-i", str(video), "-vframes", "1", "-q:v", "2", str(poster)],
        capture_output=True,
        check=True,
    )
    return poster


def _fit_in_box(vw: int, vh: int, box_l: float, box_t: float, box_w: float, box_h: float) -> tuple[float, float, float, float]:
    ar = vw / vh
    if box_w / box_h > ar:
        h = box_h
        w = h * ar
    else:
        w = box_w
        h = w / ar
    return box_l + (box_w - w) / 2, box_t + (box_h - h) / 2, w, h


def _clear_slide(slide) -> None:
    """도형 + 이전 영상용 p:timing 제거. timing을 남기면 spTgt가 쌓여 PPT 복구·재생 불가."""
    while len(slide.shapes):
        el = slide.shapes[0].element
        el.getparent().remove(el)
    sld = slide._element
    for timing in list(sld.findall(qn("p:timing"))):
        sld.remove(timing)


def _count_video_targets(slide) -> int:
    return len(slide._element.findall(".//" + qn("p:spTgt")))


def _rounded_rect(slide, left, top, w, h, fill, line=None, line_w=1.0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    return sh


def _slide_header(slide, number: str, badge: str, title: str, subtitle: str) -> None:
    """HTML .slide-header"""
    y = 0.42
    # 번호 배지
    num = _rounded_rect(slide, 0.55, y, 0.38, 0.38, RGBColor(0x10, 0xB9, 0x81), C_GREEN_LT)
    num.fill.transparency = 0.88
    ntf = num.text_frame
    ntf.text = number
    ntf.paragraphs[0].font.size = Pt(11)
    ntf.paragraphs[0].font.bold = True
    ntf.paragraphs[0].font.color.rgb = C_GREEN_LT
    ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ntf.vertical_anchor = 1

    # pill badge
    pill = _rounded_rect(slide, 1.0, y + 0.02, 1.15, 0.34, RGBColor(0x10, 0xB9, 0x81))
    pill.fill.transparency = 0.92
    pill.line.color.rgb = C_GREEN_LT
    pill.line.width = Pt(0.75)
    ptf = pill.text_frame
    ptf.text = badge.upper()
    ptf.paragraphs[0].font.size = Pt(9)
    ptf.paragraphs[0].font.bold = True
    ptf.paragraphs[0].font.color.rgb = C_GREEN_LT
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(y + 0.48), Inches(8.5), Inches(0.55))
    tf = tb.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = C_TEXT

    sb = slide.shapes.add_textbox(Inches(0.55), Inches(y + 1.02), Inches(9.0), Inches(0.4))
    sf = sb.text_frame
    sf.text = subtitle
    sf.paragraphs[0].font.size = Pt(14)
    sf.paragraphs[0].font.color.rgb = C_TEXT2


def _timeline(slide, steps: tuple[TimelineStep, ...]) -> None:
    """HTML .timeline"""
    left = 0.55
    top = 1.55
    width = 6.15
    height = 5.6

    panel = _rounded_rect(slide, left, top, width, height, C_SURFACE, C_BORDER)
    panel.fill.transparency = 0.15

    # 세로 라인
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.35), Inches(top + 0.25), Inches(0.03), Inches(height - 0.5))
    line.fill.solid()
    line.fill.fore_color.rgb = C_GREEN_LT
    line.line.fill.background()

    y = top + 0.35
    for step in steps:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.28), Inches(y + 0.08), Inches(0.14), Inches(0.14))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C_GREEN_LT
        dot.line.fill.background()

        tb = slide.shapes.add_textbox(Inches(left + 0.55), Inches(y), Inches(width - 0.7), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.text = step.title
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = C_TEXT
        p2 = tf.add_paragraph()
        p2.text = step.desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT2
        p2.space_before = Pt(2)
        y += 1.02


def _phone_frame(slide, phone_left: float) -> tuple[float, float, float, float]:
    """폰 베젤 + 검은 화면 영역 좌표 반환 (영상 삽입용)."""
    outer = _rounded_rect(
        slide,
        phone_left,
        PHONE_TOP,
        PHONE_W,
        PHONE_H,
        C_SURFACE,
        RGBColor(0xFF, 0xFF, 0xFF),
        1.5,
    )
    outer.fill.transparency = 0.05

    # 노치
    notch = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(phone_left + PHONE_W / 2 - 0.28),
        Inches(PHONE_TOP + 0.06),
        Inches(0.56),
        Inches(0.07),
    )
    notch.fill.solid()
    notch.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    notch.fill.transparency = 0.85
    notch.line.fill.background()

    sx = phone_left + SCREEN_INSET_X
    sy = PHONE_TOP + SCREEN_INSET_Y
    _rounded_rect(slide, sx, sy, SCREEN_W, SCREEN_H, RGBColor(0x00, 0x00, 0x00))
    return sx, sy, SCREEN_W, SCREEN_H


def _fullscreen_button(slide, phone_left: float, video: Path) -> float:
    """전체화면 버튼 — 슬라이드쇼에서 클릭 시 MP4를 기본 플레이어로 연다."""
    btn_w = 2.05
    btn_h = 0.36
    btn_l = phone_left + (PHONE_W - btn_w) / 2
    btn_t = PHONE_TOP + PHONE_H + 0.1
    btn = _rounded_rect(slide, btn_l, btn_t, btn_w, btn_h, C_GREEN_LT, C_GREEN, 1.25)
    tf = btn.text_frame
    tf.text = "⛶ 전체화면으로 보기"
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = 1
    btn.click_action.hyperlink.address = video.resolve().as_uri()
    return btn_t + btn_h


def _play_hint(slide, phone_left: float, top_after_btn: float, extra: str = "") -> None:
    hint = slide.shapes.add_textbox(
        Inches(phone_left - 0.15),
        Inches(top_after_btn + 0.08),
        Inches(PHONE_W + 0.3),
        Inches(0.42),
    )
    tf = hint.text_frame
    tf.text = "▶ F5 슬라이드쇼: 영상 클릭 재생 · 위 버튼: 전체화면(외부 플레이어)" + extra
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = C_MUTED
    p.alignment = PP_ALIGN.CENTER


def _embed_video(slide, video: Path, screen_l: float, screen_t: float, screen_w: float, screen_h: float) -> None:
    """포스터를 add_movie에 넘겨 스피커 아이콘 대신 첫 프레임(또는 지정 PNG) 표시."""
    vw, vh = _probe_video_size(video)
    left, top, w, h = _fit_in_box(vw, vh, screen_l, screen_t, screen_w, screen_h)
    poster = _poster_for_video(video)
    slide.shapes.add_movie(
        str(video.resolve()),
        Inches(left),
        Inches(top),
        width=Inches(w),
        height=Inches(h),
        poster_frame_image=str(poster.resolve()),
        mime_type="video/mp4",
    )


def _demo_narration_note(slide) -> None:
    """시연 영상 나레이션 — Hugging Face 로컬 TTS 안내."""
    box = slide.shapes.add_textbox(Inches(0.55), Inches(6.38), Inches(12.2), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🎙️ 나레이션: Hugging Face에서 받은 Supertonic 3 로컬 TTS로 합성 (무료·오프라인)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = C_TEXT2
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = HF_TTS_MODEL_URL
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_GREEN_LT
    p2.alignment = PP_ALIGN.CENTER
    try:
        p2.hyperlink.address = HF_TTS_MODEL_URL
    except AttributeError:
        pass


def _build_slide_12(slide) -> None:
    _slide_header(slide, "11", "Demo", "데모 시나리오", "발표 시 실시간 연출")
    _timeline(slide, DEMO_STEPS)
    screen = _phone_frame(slide, PHONE_LEFT_DEMO)
    demo = ASSETS / "recycle_demo.mp4"
    y = _fullscreen_button(slide, PHONE_LEFT_DEMO, demo)
    _play_hint(slide, PHONE_LEFT_DEMO, y)
    _embed_video(slide, demo, *screen)
    _demo_narration_note(slide)


def _ad_production_note(slide) -> None:
    """SNS 광고 영상 — Cursor Higgsfield MCP 안내."""
    box = slide.shapes.add_textbox(Inches(1.2), Inches(6.35), Inches(10.9), Inches(0.95))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "🎬 제작: Cursor의 Higgsfield MCP로 AI 홍보 영상 생성 (9:16 UGC · 한국어 립싱크)"
    p1.font.size = Pt(11)
    p1.font.color.rgb = C_TEXT2
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"참고: {HIGGSFIELD_URL} · docs/presentation/HIGGSFIELD_AD_VIDEO.md"
    p2.font.size = Pt(10)
    p2.font.color.rgb = C_GREEN_LT
    p2.alignment = PP_ALIGN.CENTER
    try:
        p2.hyperlink.address = HIGGSFIELD_URL
    except AttributeError:
        pass


def _build_slide_25(slide) -> None:
    _slide_header(
        slide,
        "25",
        "Marketing Video",
        "SNS 홍보 영상 (15초)",
        "릴스·틱톡·쇼츠 · 9:16 UGC",
    )
    screen = _phone_frame(slide, PHONE_LEFT_AD)
    ad = ASSETS / "recycle_ad.mp4"
    y = _fullscreen_button(slide, PHONE_LEFT_AD, ad)
    _play_hint(slide, PHONE_LEFT_AD, y, " · 음소거 해제 권장")
    _embed_video(slide, ad, *screen)
    _ad_production_note(slide)

    foot = slide.shapes.add_textbox(Inches(2.5), Inches(7.05), Inches(8.3), Inches(0.28))
    foot.text_frame.text = "assets/recycle_ad.mp4"
    foot.text_frame.paragraphs[0].font.size = Pt(9)
    foot.text_frame.paragraphs[0].font.color.rgb = C_MUTED
    foot.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def rebuild_video_slides(ppt_path: Path) -> Path:
    base = ROOT / "RecycleAI_Project_Deck.pptx"
    if not ppt_path.is_file():
        ppt_path = base if base.is_file() else ROOT / "RecycleAI_Project_Deck_with_videos.pptx"
    prs = Presentation(str(ppt_path))
    n = len(prs.slides)
    for page in (12, 25):
        if page > n:
            raise ValueError(f"슬라이드 {page} 없음 (총 {n}장)")

    builders = {12: _build_slide_12, 25: _build_slide_25}
    for page, builder in builders.items():
        slide = prs.slides[page - 1]
        _clear_slide(slide)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = C_BG
        builder(slide)
        targets_n = _count_video_targets(slide)
        if targets_n != 1:
            raise RuntimeError(f"슬라이드 {page}: 영상 타깃 {targets_n}개(1개여야 함). PPT 손상 방지를 위해 중단.")

    out = ROOT / "RecycleAI_Project_Deck_with_videos.pptx"
    targets = [out]
    if ppt_path.resolve() != out.resolve():
        targets.append(ppt_path)
    saved = None
    for target in targets:
        try:
            prs.save(str(target))
            saved = target
            break
        except PermissionError:
            continue
    if saved is None:
        alt = ROOT / "RecycleAI_Project_Deck_videos_redesign.pptx"
        prs.save(str(alt))
        saved = alt
        print(f"파일이 열려 있어 저장: {alt}")
    return saved


def embed_videos(ppt_path: Path | None = None, **_kwargs) -> Path:
    path = ppt_path or ROOT / "RecycleAI_Project_Deck_with_videos.pptx"
    return rebuild_video_slides(path)


if __name__ == "__main__":
    src = ROOT / "RecycleAI_Project_Deck.pptx"
    if not src.is_file():
        src = ROOT / "RecycleAI_Project_Deck_with_videos.pptx"
    path = rebuild_video_slides(src)
    print(f"12·25페이지 재디자인 완료: {path}")
